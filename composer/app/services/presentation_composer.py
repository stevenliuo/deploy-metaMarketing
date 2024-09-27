import io
from dataclasses import dataclass

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.slide import Slide
from pptx.util import Inches, Pt

class SlideDoesntExist(Exception):
    """ Raised when no such slide exists. """

    def __init__(self):
        super().__init__(f"Slide doesn't exist")


@dataclass
class SlideData:
    title: str
    content: str
    image: bytes | None
    option: int


class PresentationComposer:
    def __init__(self):
        super().__init__()

    def compose(self, template: bytes,
                title: str, subtitle: str | None, footer: str | None,
                slides: list[SlideData]) -> bytes:
        input_file = io.BytesIO(template)
        prs: Presentation = Presentation(input_file)

        first_slide = prs.slides[0]

        title_placeholder, subtitle_placeholder = self._find_title_and_subtitle_shapes(first_slide)
        if title_placeholder:
            new_text = title or ''
            self._update_text_with_formatting2(title_placeholder, new_text)
        if subtitle_placeholder and subtitle:
            self._update_text_with_formatting2(subtitle_placeholder, subtitle)

        for i, slide_data in enumerate(slides):
            slide_no = i + 1
            if slide_no > len(prs.slides):
                raise SlideDoesntExist()
            slide = prs.slides[slide_no]

            # Add title
            title_placeholder, subtitle_placeholder = self._find_title_and_subtitle_shapes(slide)
            title_bottom = None
            if slide_data.option == 4:
                title_bottom = 0
            else:
                if title_placeholder:
                    new_text = slide_data.title or ''
                    self._update_text_with_formatting2(title_placeholder, new_text)
                    title_bottom = title_placeholder.top + title_placeholder.height
                else:
                    title_shape = self._find_first_textbox_shape(slide)
                    new_text = slide_data.title or ''
                    if title_shape:
                        self._update_text_with_formatting2(title_shape, new_text)
                        title_bottom = title_shape.top + title_shape.height

            text_content = slide_data.content.replace('**', '')

            if slide_data.image is None:
                self._add_full_slide_text(prs, slide, text_content, title_bottom=title_bottom)
            else:
                if slide_data.option == 1:
                    self._add_left_half_slide_text(prs, slide, text_content, title_bottom=title_bottom)
                    self._add_right_half_slide_image2(prs, slide, slide_data.image, title_bottom=title_bottom)
                if slide_data.option == 3 or slide_data.option == 4:
                    self._add_full_slide_image(prs, slide, slide_data.image, title_bottom=title_bottom)
                if slide_data.option == 5:
                    text_height = self._add_half_height_slide_text(prs, slide, text_content, title_bottom=title_bottom)
                    self._add_full_slide_image(prs, slide, slide_data.image, title_bottom=int(title_bottom + text_height))
            # Add footer
            if footer:
                footer_placeholder = self._find_footer_shape(slide)
                if footer_placeholder:
                    self._update_text_with_formatting2(footer_placeholder, footer)
                else:
                    footer_placeholder = self._find_footer_shape(prs.slide_master)
                    if footer_placeholder:
                        slide.shapes.clone_placeholder(footer_placeholder)
                        footer_placeholder = self._find_footer_shape(slide)
                        self._update_text_with_formatting2(footer_placeholder, footer)

        # Save the modified presentation
        output_file = io.BytesIO()
        prs.save(output_file)
        return output_file.getvalue()

    def _find_title_and_subtitle_shapes(self, slide: Slide):
        title = None
        subtitle = None
        for placeholder in slide.placeholders:
            if title is None and (placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE or
                                  placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.CENTER_TITLE):
                title = placeholder
            elif subtitle is None and (placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE):
                subtitle = placeholder

        return title, subtitle

    def _find_footer_shape(self, slide: Slide):
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == PP_PLACEHOLDER_TYPE.FOOTER:
                return placeholder
        return None

    def _find_first_textbox_shape(self, slide: Slide):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return shape
        return None

    def _update_text_with_formatting2(self, shape, new_text: str):
        try:
            paragraph = shape.text_frame.paragraphs[0]
            # Delete all runs except the first one
            while len(paragraph.runs) > 1:
                paragraph.runs[-1].text = ''
                del paragraph.runs[-1]
            if not len(paragraph.runs):
                paragraph.add_run()
            paragraph.runs[0].text = new_text
        except Exception as e:
            shape.text = new_text

    def _calculate_text_top_and_height(self, slide_height, title_bottom, default_top=Inches(0.5)):
        if title_bottom is None or title_bottom > Inches(2):
            top = default_top
            height = slide_height - Inches(1)
        else:
            top = title_bottom
            height = slide_height - top
        return top, height

    def _add_full_slide_text(self, prs: Presentation, slide, text: str, *, title_bottom=None, margin=Inches(0.5)):
        left = margin
        width = prs.slide_width - 2 * margin
        top, height = self._calculate_text_top_and_height(prs.slide_height, title_bottom)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        p = text_frame.add_paragraph()

        if len(text) < 3600:
            p.font.size = Pt(10)
        else:
            p.font.size = Pt(8)
        p.text = text

        return height
    
    def _add_half_height_slide_text(self, prs: Presentation, slide, text: str, *, title_bottom=None, margin=Inches(0.5)):
        left = margin
        width = prs.slide_width - 2 * margin
        top, height = self._calculate_text_top_and_height(prs.slide_height, title_bottom)
        height /= 4

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        p = text_frame.add_paragraph()

        if len(text) < 3600:
            p.font.size = Pt(10)
        else:
            p.font.size = Pt(8)
        p.text = text

        return height

    def _add_left_half_slide_text(self, prs: Presentation, slide, text: str, *, title_bottom=None):
        left = Inches(0.5)
        width = int(prs.slide_width / 2 - Inches(1))
        top, height = self._calculate_text_top_and_height(prs.slide_height, title_bottom)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        p = text_frame.add_paragraph()
        if len(text) < 900:
            p.font.size = Pt(12)
        elif len(text) < 1800:
            p.font.size = Pt(10)
        elif len(text) < 3600:
            p.font.size = Pt(8)
        else:
            p.font.size = Pt(6)

        p.text = text

    def _add_right_half_slide_image(self, slide, image: bytes):
        left = Inches(4.5)
        top = Inches(1)
        height = Inches(4)
        image_file = io.BytesIO(image)
        slide.shapes.add_picture(image_file=image_file, left=left, top=top, height=height)

        img = Image.open(io.BytesIO(image))
        img_width, img_height = img.size

        # Calculate the scaling factor to fit within the specified dimensions
        scaling_factor = min(4 / img_width, 4 / img_height)
        new_width = int(img_width * scaling_factor)
        new_height = int(img_height * scaling_factor)

        # Add the image to the slide
        pic = slide.shapes.add_picture(io.BytesIO(image), left, top, width=Inches(new_width / 96), height=Inches(new_height / 96))

    def _add_right_half_slide_image2(self, prs: Presentation, slide: Slide, image: bytes, title_bottom=0, margin=Inches(0.5)):
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = int(prs.slide_height - title_bottom)

        img = Image.open(io.BytesIO(image))
        img_width, img_height = img.size
        img_ratio = img_width / img_height
        half_slide_ratio = slide_width / 2 / slide_height

        # Calculate image position and size
        if img_ratio > half_slide_ratio:
            # Image is wider than half slide
            width = int(slide_width / 2 - 2 * margin)
            height = int(width / img_ratio)
            left = int(slide_width / 2)
            top = int(title_bottom + margin)
        else:
            # Image is taller than half slide
            height = int(slide_height - 2 * margin)
            width = int(height * img_ratio)
            left = int(slide_width / 2 + (slide_width / 2 - width) / 2)
            top = int(title_bottom + margin)

        # Add picture to slide
        pic = slide.shapes.add_picture(io.BytesIO(image), left, top, width, height)

        # Adjust picture position to fit in right half
        pic.left = left
        pic.top = top

    def _add_full_slide_image(self, prs: Presentation, slide: Slide, image: bytes, title_bottom=0, margin=Inches(0.5)):
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = int(prs.slide_height - title_bottom)

        img = Image.open(io.BytesIO(image))
        img_width, img_height = img.size
        img_ratio = img_width / img_height
        slide_ratio = slide_width / slide_height

        # Calculate image position and size
        if img_ratio > slide_ratio:
            # Image is wider than half slide
            width = int(slide_width - 2 * margin)
            height = int(width / img_ratio)
            left = margin
            top = int(title_bottom + margin)
        else:
            # Image is taller than half slide
            height = int(slide_height - 2 * margin)
            width = int(height * img_ratio)
            left = int((slide_width - width) / 2)
            top = int(title_bottom + margin)

        # Add picture to slide
        slide.shapes.add_picture(io.BytesIO(image), left, top, width, height)